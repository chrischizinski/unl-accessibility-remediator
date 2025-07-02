"""
PowerPoint Processor Module

Handles .pptx files for accessibility analysis and remediation.
Extracts content, analyzes accessibility issues, and applies fixes following UNL's
WCAG 2.1 Level AA requirements and Title II compliance guidelines.
"""

import logging
import time
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.shapes.group import GroupShape
from pptx.text.text import TextFrame
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .ai_assistant import AIAssistant, AccessibilityResults, SlideAnalysis
from .contrast_checker import ContrastChecker
from .alt_text_generator import AltTextGenerator
from .link_checker import LinkChecker


class PowerPointProcessor:
    """Processes PowerPoint presentations for accessibility compliance."""
    
    def __init__(self, ai_assistant: AIAssistant):
        """
        Initialize PowerPoint processor with AI assistant.
        
        Args:
            ai_assistant: AI assistant for accessibility analysis
        """
        self.ai_assistant = ai_assistant
        self.contrast_checker = ContrastChecker()
        self.alt_text_generator = AltTextGenerator(ai_assistant)
        self.link_checker = LinkChecker(ai_assistant)
        self.logger = logging.getLogger(__name__)
    
    def analyze_accessibility(self, pptx_path: Path) -> AccessibilityResults:
        """
        Analyze PowerPoint presentation for accessibility issues.
        
        Args:
            pptx_path: Path to .pptx file
            
        Returns:
            AccessibilityResults with analysis for all slides
        """
        start_time = time.time()
        self.logger.info(f"Starting accessibility analysis of {pptx_path}")
        
        try:
            presentation = Presentation(str(pptx_path))
            slide_analyses = []
            
            for i, slide in enumerate(presentation.slides, 1):
                self.logger.debug(f"Analyzing slide {i}/{len(presentation.slides)}")
                
                # Extract slide data
                slide_data = self._extract_slide_data(slide, i)
                
                # Get AI analysis
                analysis = self.ai_assistant.analyze_slide(slide_data)
                analysis.title = slide_data.get('title')
                
                slide_analyses.append(analysis)
            
            # Calculate overall metrics
            total_issues = sum(
                len(s.alt_text_suggestions) + len(s.link_improvements) + 
                len(s.contrast_issues) + len(s.content_issues) 
                for s in slide_analyses
            )
            
            auto_fixable_count = sum(len(s.auto_fixable) for s in slide_analyses)
            manual_review_count = sum(len(s.manual_review) for s in slide_analyses)
            
            # Calculate overall accessibility score (0-100)
            if total_issues == 0:
                overall_score = 100.0
            else:
                # Weight by confidence and severity
                weighted_score = sum(s.confidence_score for s in slide_analyses) / len(slide_analyses)
                issue_penalty = min(total_issues * 5, 80)  # Cap penalty at 80 points
                overall_score = max(20.0, weighted_score * 100 - issue_penalty)
            
            processing_time = time.time() - start_time
            
            results = AccessibilityResults(
                slides=slide_analyses,
                overall_score=overall_score,
                total_issues=total_issues,
                auto_fixable_count=auto_fixable_count,
                manual_review_count=manual_review_count,
                processing_time=processing_time
            )
            
            self.logger.info(f"Analysis completed: {total_issues} issues found, "
                           f"score: {overall_score:.1f}/100")
            
            return results
            
        except Exception as e:
            self.logger.error(f"Failed to analyze PowerPoint: {e}")
            raise
    
    def apply_fixes(self, results: AccessibilityResults, input_path: Path, 
                   output_dir: Path) -> AccessibilityResults:
        """
        Apply automatic fixes to PowerPoint presentation.
        
        Args:
            results: Analysis results
            input_path: Original .pptx file path
            output_dir: Directory for output files
            
        Returns:
            Updated AccessibilityResults with applied fixes
        """
        self.logger.info("Applying automatic accessibility fixes")
        
        try:
            presentation = Presentation(str(input_path))
            applied_fixes = []
            
            for slide_idx, (slide, analysis) in enumerate(zip(presentation.slides, results.slides)):
                # Apply alt text fixes
                alt_fixes = self._apply_alt_text_fixes(slide, analysis)
                applied_fixes.extend(alt_fixes)
                
                # Apply link text improvements
                link_fixes = self._apply_link_fixes(slide, analysis)
                applied_fixes.extend(link_fixes)
                
                # Apply title improvements if available
                title_fixes = self._apply_title_fixes(slide, analysis)
                applied_fixes.extend(title_fixes)
            
            # Save modified presentation
            output_file = output_dir / f"{input_path.stem}_accessible.pptx"
            presentation.save(str(output_file))
            
            self.logger.info(f"Applied {len(applied_fixes)} fixes, saved to {output_file}")
            
            # Update results with applied fixes
            results.automatic_fixes = applied_fixes
            
            return results
            
        except Exception as e:
            self.logger.error(f"Failed to apply fixes: {e}")
            raise
    
    def _extract_slide_data(self, slide, slide_number: int) -> Dict[str, Any]:
        """Extract structured data from a PowerPoint slide."""
        data = {
            'slide_number': slide_number,
            'title': None,
            'text_content': [],
            'images': [],
            'links': [],
            'colors': []
        }
        
        # Extract content from all shapes
        for shape in slide.shapes:
            self._process_shape(shape, data)
        
        # Try to identify slide title
        if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
            data['title'] = slide.shapes.title.text.strip()
        elif data['text_content']:
            # Use first text block as potential title
            data['title'] = data['text_content'][0][:100] + "..." if len(data['text_content'][0]) > 100 else data['text_content'][0]
        
        return data
    
    def _process_shape(self, shape: BaseShape, data: Dict[str, Any]) -> None:
        """Process a single shape and extract relevant accessibility data."""
        try:
            # Handle text content
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text_content = self._extract_text_content(shape.text_frame)
                if text_content.strip():
                    data['text_content'].append(text_content)
                    
                    # Extract color information
                    colors = self._extract_text_colors(shape.text_frame)
                    data['colors'].extend(colors)
                    
                    # Extract hyperlinks
                    links = self._extract_hyperlinks(shape.text_frame)
                    data['links'].extend(links)
            
            # Handle images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_data = self._extract_image_data(shape)
                if image_data:
                    data['images'].append(image_data)
            
            # Handle grouped shapes recursively
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for grouped_shape in shape.shapes:
                    self._process_shape(grouped_shape, data)
                    
        except Exception as e:
            self.logger.debug(f"Error processing shape: {e}")
    
    def _extract_text_content(self, text_frame: TextFrame) -> str:
        """Extract all text content from a text frame."""
        text_parts = []
        for paragraph in text_frame.paragraphs:
            para_text = ""
            for run in paragraph.runs:
                para_text += run.text
            if para_text.strip():
                text_parts.append(para_text.strip())
        return "\n".join(text_parts)
    
    def _extract_text_colors(self, text_frame: TextFrame) -> List[Dict[str, Any]]:
        """Extract text color information for contrast checking."""
        colors = []
        
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.color and run.font.color.rgb:
                        # Get text color
                        text_color = run.font.color.rgb
                        
                        # For background, we'd need slide background or shape fill
                        # This is a simplified version - full implementation would
                        # need to traverse the shape hierarchy for accurate background color
                        
                        colors.append({
                            'element_type': 'text',
                            'foreground': f'#{text_color.r:02x}{text_color.g:02x}{text_color.b:02x}',
                            'background': '#ffffff',  # Assume white background - TODO: improve
                            'text': run.text[:50] + "..." if len(run.text) > 50 else run.text
                        })
        except Exception as e:
            self.logger.debug(f"Error extracting colors: {e}")
        
        return colors
    
    def _extract_hyperlinks(self, text_frame: TextFrame) -> List[Dict[str, str]]:
        """Extract hyperlink information."""
        links = []
        
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.hyperlink and run.hyperlink.address:
                        links.append({
                            'text': run.text,
                            'url': run.hyperlink.address
                        })
        except Exception as e:
            self.logger.debug(f"Error extracting hyperlinks: {e}")
        
        return links
    
    def _extract_image_data(self, shape: Picture) -> Optional[Dict[str, Any]]:
        """Extract image information for alt text analysis."""
        try:
            # Get existing alt text if present
            alt_text = getattr(shape, 'alt_text', '') or ''
            
            # Create image identifier
            image_id = f"img_{id(shape)}"
            
            # Get basic image info
            image_data = {
                'id': image_id,
                'alt_text': alt_text,
                'description': f"Image with dimensions {shape.width} x {shape.height}",
                'has_alt_text': bool(alt_text.strip()),
                'shape_ref': shape  # Keep reference for applying fixes
            }
            
            return image_data
            
        except Exception as e:
            self.logger.debug(f"Error extracting image data: {e}")
            return None
    
    def _apply_alt_text_fixes(self, slide, analysis: SlideAnalysis) -> List[Dict[str, Any]]:
        """Apply alt text fixes to slide images."""
        fixes = []
        
        for suggestion in analysis.alt_text_suggestions:
            image_id = suggestion.get('image_id')
            suggested_alt = suggestion.get('suggested_alt', '')
            
            if not suggested_alt:
                continue
            
            # Find the image shape by ID (simplified lookup)
            for shape in slide.shapes:
                if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE and 
                    f"img_{id(shape)}" == image_id):
                    
                    try:
                        # Apply alt text
                        shape.alt_text = suggested_alt
                        
                        fixes.append({
                            'type': 'alt_text',
                            'slide_number': analysis.slide_number,
                            'description': f"Added alt text: '{suggested_alt[:50]}...'",
                            'element': image_id
                        })
                        
                        self.logger.debug(f"Applied alt text to {image_id}")
                        break
                        
                    except Exception as e:
                        self.logger.warning(f"Failed to apply alt text to {image_id}: {e}")
        
        return fixes
    
    def _apply_link_fixes(self, slide, analysis: SlideAnalysis) -> List[Dict[str, Any]]:
        """Apply link text improvements to slide."""
        fixes = []
        
        for improvement in analysis.link_improvements:
            original_text = improvement.get('original_text', '')
            suggested_text = improvement.get('suggested_text', '')
            
            if not original_text or not suggested_text:
                continue
            
            # Find and update link text in slide shapes
            updated = self._update_link_text(slide, original_text, suggested_text)
            
            if updated:
                fixes.append({
                    'type': 'link_text',
                    'slide_number': analysis.slide_number,
                    'description': f"Updated link text: '{original_text}' → '{suggested_text}'",
                    'element': 'hyperlink'
                })
        
        return fixes
    
    def _apply_title_fixes(self, slide, analysis: SlideAnalysis) -> List[Dict[str, Any]]:
        """Apply slide title improvements."""
        fixes = []
        
        if analysis.suggested_title and slide.shapes.title:
            try:
                old_title = slide.shapes.title.text
                slide.shapes.title.text = analysis.suggested_title
                
                fixes.append({
                    'type': 'title',
                    'slide_number': analysis.slide_number,
                    'description': f"Updated title: '{old_title}' → '{analysis.suggested_title}'",
                    'element': 'slide_title'
                })
                
            except Exception as e:
                self.logger.warning(f"Failed to update slide title: {e}")
        
        return fixes
    
    def _update_link_text(self, slide, original_text: str, new_text: str) -> bool:
        """Update hyperlink text in slide shapes."""
        updated = False
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if (run.text.strip() == original_text.strip() and 
                            run.hyperlink and run.hyperlink.address):
                            
                            try:
                                run.text = new_text
                                updated = True
                                self.logger.debug(f"Updated link text: {original_text} → {new_text}")
                            except Exception as e:
                                self.logger.warning(f"Failed to update link text: {e}")
        
        return updated